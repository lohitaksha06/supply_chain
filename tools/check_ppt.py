from pptx import Presentation

prs = Presentation('PharmaChain_Final_Deliverable.pptx')
print('slides', len(prs.slides))
for i, slide in enumerate(prs.slides, start=1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            t = (shape.text or '').strip()
            if t:
                texts.append(t.replace('\n', ' / '))
    print(f"{i:02d}: {len(texts)} text blocks")
    for t in texts[:2]:
        print('  -', t[:160])
