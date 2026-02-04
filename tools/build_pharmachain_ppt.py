from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "ppt.pptx"
OUT_DIR = ROOT / "generated"
OUT_PPTX = ROOT / "PharmaChain_Final_Deliverable.pptx"


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        ("arialbd.ttf" if bold else "arial.ttf"),
        ("calibrib.ttf" if bold else "calibri.ttf"),
        ("DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"),
    ]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _wrap_lines(text: str, max_chars: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current: list[str] = []
    for w in words:
        if sum(len(x) for x in current) + len(current) + len(w) > max_chars and current:
            lines.append(" ".join(current))
            current = [w]
        else:
            current.append(w)
    if current:
        lines.append(" ".join(current))
    return lines


def draw_decision_matrix(path: Path) -> None:
    w, h = 1800, 1000
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)

    title_font = _font(44, bold=True)
    header_font = _font(26, bold=True)
    cell_font = _font(22)

    x0, y0 = 80, 120
    table_w, table_h = w - 160, h - 200
    cols = [
        ("Opportunity", 0.30),
        ("Impact", 0.12),
        ("Feasibility", 0.14),
        ("Regulatory fit", 0.16),
        ("Blockchain fit", 0.14),
        ("Total", 0.14),
    ]
    rows = [
        ("Pharma supply chain traceability", [5, 4, 5, 5]),
        ("Food supply chain traceability", [4, 4, 4, 4]),
        ("Luxury goods anti-counterfeit", [3, 4, 2, 4]),
        ("Academic credential verification", [3, 5, 2, 4]),
    ]

    d.text((80, 40), "Decision Matrix", fill=(40, 40, 40), font=title_font)
    d.text((420, 52), "(scores 1–5)", fill=(90, 90, 90), font=_font(22))

    col_x = [x0]
    for _, frac in cols:
        col_x.append(col_x[-1] + int(table_w * frac))

    header_h = 64
    row_h = int((table_h - header_h) / (len(rows) + 0.2))

    d.rounded_rectangle([x0, y0, x0 + table_w, y0 + header_h], radius=12, fill=(19, 131, 198))

    for i, (name, _) in enumerate(cols):
        x_left, x_right = col_x[i], col_x[i + 1]
        d.line([x_right, y0, x_right, y0 + table_h], fill=(220, 220, 220), width=2)
        d.text((x_left + 14, y0 + 18), name, fill="white", font=header_font)

    y = y0 + header_h
    for r_idx, (opp, scores) in enumerate(rows):
        total = sum(scores)
        is_highlight = r_idx == 0
        bg = (230, 247, 255) if is_highlight else (255, 255, 255)
        d.rectangle([x0, y, x0 + table_w, y + row_h], fill=bg)
        d.line([x0, y + row_h, x0 + table_w, y + row_h], fill=(230, 230, 230), width=2)

        # Wrap opportunity
        opp_lines = _wrap_lines(opp, max_chars=34)
        yy = y + 10
        for ln in opp_lines[:2]:
            d.text((col_x[0] + 14, yy), ln, fill=(30, 30, 30), font=cell_font)
            yy += 28

        for c, val in enumerate(scores, start=1):
            d.text((col_x[c] + 24, y + 16), str(val), fill=(30, 30, 30), font=cell_font)
        d.text((col_x[5] + 24, y + 16), str(total), fill=(30, 30, 30), font=cell_font)

        if is_highlight:
            d.text((x0 + table_w - 260, y + 16), "Selected", fill=(19, 131, 198), font=_font(22, bold=True))

        y += row_h

    d.rounded_rectangle([x0, y0, x0 + table_w, y0 + table_h], radius=12, outline=(200, 200, 200), width=3)
    img.save(path)


def draw_bmc(path: Path) -> None:
    w, h = 1800, 1000
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)

    title_font = _font(44, bold=True)
    h_font = _font(24, bold=True)
    b_font = _font(20)

    d.text((80, 40), "Business Model Canvas — PharmaChain", fill=(40, 40, 40), font=title_font)

    x0, y0 = 80, 120
    canvas_w, canvas_h = w - 160, h - 200

    col_fracs = [0.18, 0.18, 0.18, 0.28, 0.18]
    row_fracs = [0.62, 0.38]

    xs = [x0]
    for f in col_fracs:
        xs.append(xs[-1] + int(canvas_w * f))
    ys = [y0, y0 + int(canvas_h * row_fracs[0]), y0 + canvas_h]

    def box(x1: int, y1: int, x2: int, y2: int, title: str, bullets: list[str]) -> None:
        d.rounded_rectangle([x1, y1, x2, y2], radius=12, outline=(200, 200, 200), width=3)
        d.rectangle([x1, y1, x2, y1 + 46], fill=(230, 247, 255))
        d.text((x1 + 12, y1 + 10), title, fill=(19, 131, 198), font=h_font)
        y = y1 + 60
        max_chars = max(18, int((x2 - x1) / 13))
        for b in bullets:
            text = "• " + b
            lines = _wrap_lines(text, max_chars=max_chars)
            for line in lines[:3]:
                d.text((x1 + 14, y), line, fill=(40, 40, 40), font=b_font)
                y += 26
            y += 6

    key_partners = [
        "Manufacturers, distributors, wholesalers",
        "Retail pharmacies & hospitals",
        "Regulators (track & trace compliance)",
        "GS1/serialization ecosystem",
        "Cloud + security/audit partners",
    ]
    key_activities = [
        "Onboard participants & integrate ERPs",
        "Operate permissioned blockchain nodes",
        "Maintain smart contracts & APIs",
        "Data validation + compliance reporting",
        "Security monitoring & audits",
    ]
    key_resources = [
        "Rust backend API + database",
        "Next.js web app (role-based dashboards)",
        "Merkle proofs + RSA signatures",
        "Blockchain governance + node infra",
        "Analytics + audit trail storage",
    ]
    value_props = [
        "End-to-end batch/pack provenance",
        "Anti-counterfeit verification (QR scan)",
        "Faster targeted recalls",
        "Shared source of truth",
        "Lower reconciliation costs",
    ]
    cust_relationships = [
        "Enterprise onboarding & support",
        "Compliance/audit assistance",
        "Dashboards + alerts",
        "Developer APIs + docs",
    ]
    channels = [
        "Direct enterprise sales",
        "Partnerships with logistics/IT vendors",
        "Web portal for stakeholders",
        "Patient verification via QR/mobile web",
    ]
    cust_segments = [
        "Manufacturers (brand protection)",
        "Distributors (custody tracking)",
        "Pharmacies/hospitals (dispensing)",
        "Patients (authenticity check)",
        "Regulators (audit)",
    ]
    cost_structure = [
        "Cloud hosting + node ops",
        "Development + security audits",
        "Integrations + customer success",
        "Support, compliance, governance",
    ]
    revenue_streams = [
        "SaaS subscription per participant",
        "Tiered pricing by volume",
        "Premium analytics + compliance reports",
        "Integration/pro services",
    ]

    box(xs[0], ys[0], xs[1], ys[1], "Key Partners", key_partners)
    box(xs[1], ys[0], xs[2], ys[1], "Key Activities", key_activities)
    box(xs[2], ys[0], xs[3], ys[1], "Key Resources", key_resources)
    box(xs[3], ys[0], xs[4], ys[1], "Value Propositions", value_props)
    box(xs[4], ys[0], xs[5], ys[1], "Customer Segments", cust_segments)

    mid_y = ys[0] + int((ys[1] - ys[0]) * 0.62)
    box(xs[3], mid_y, xs[4], ys[1], "Customer Relationships", cust_relationships)
    box(xs[4], mid_y, xs[5], ys[1], "Channels", channels)

    box(xs[0], ys[1], xs[3], ys[2], "Cost Structure", cost_structure)
    box(xs[3], ys[1], xs[5], ys[2], "Revenue Streams", revenue_streams)

    img.save(path)


def _set_text(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    # Keep it simple: new lines become paragraphs.
    lines = [ln.strip() for ln in text.split("\n")]
    if not lines:
        return
    tf.paragraphs[0].text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln


def main() -> None:
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    OUT_DIR.mkdir(exist_ok=True)
    decision_img = OUT_DIR / "decision_matrix.png"
    bmc_img = OUT_DIR / "business_model_canvas.png"

    draw_decision_matrix(decision_img)
    draw_bmc(bmc_img)

    prs = Presentation(str(TEMPLATE))

    # Slide 1 — Title Slide
    s1 = prs.slides[0]
    _set_text(s1.placeholders[0], "PharmaChain")
    _set_text(s1.placeholders[1], "Lohitaksha Patary")
    _set_text(s1.placeholders[11], "Undergraduate Student, Amrita Vishwa Vidyapeetham, Bangalore")
    _set_text(s1.placeholders[10], "02/04/2026")

    # Slide 2 — Industry Analysis
    s2 = prs.slides[1]
    _set_text(s2.placeholders[13], "Pharmaceutical supply chain (Healthcare)")
    _set_text(s2.placeholders[14], "Hospitals, pharmacies, and patients needing verifiable batch-level provenance")
    _set_text(
        s2.placeholders[15],
        "• High counterfeit risk and patient safety impact\n"
        "• Many parties (manufacturer → distributor → pharmacy/hospital → customer) need a shared source of truth\n"
        "• Regulators require traceability and fast, targeted recalls\n"
        "• Current ERP/track-and-trace data is siloed and hard to reconcile across organizations\n"
        "• Blockchain provides tamper-evident audit trails, provenance, and cryptographic verification (hash chain + signatures)\n"
        "• Merkle proofs enable efficient integrity checks without exposing all data",
    )

    # Slide 3 — Market Research
    s3 = prs.slides[2]
    _set_text(
        s3.placeholders[16],
        "Approach:\n"
        "• Desk research on pharma counterfeits, recalls, and serialization/track-and-trace regulations\n"
        "• Stakeholder mapping (manufacturers, distributors, hospitals, pharmacies, customers, regulators)\n"
        "• Process analysis of batch creation, custody transfer, and dispensing events\n"
        "• Gap analysis of existing systems (ERP + EDI + centralized portals)\n\n"
        "Key findings:\n"
        "• Data fragmentation creates blind spots and disputes during recalls/investigations\n"
        "• Last-mile verification for customers is weak (hard to prove authenticity)\n"
        "• A permissioned, shared ledger with digital signatures improves trust and accountability\n"
        "• Storing proofs (hashes, Merkle roots) gives integrity guarantees with minimal data sharing",
    )

    # Slide 4 — Competitive Analysis
    s4 = prs.slides[3]
    _set_text(
        s4.placeholders[16],
        "Top competitors (examples):\n"
        "1) MediLedger (Chronicled)\n"
        "   Strengths: Industry consortia, DSCSA alignment, strong network effects\n"
        "   Weaknesses: Adoption/governance complexity; integration effort\n"
        "2) TraceLink\n"
        "   Strengths: Established pharma network; robust serialization & compliance tooling\n"
        "   Weaknesses: Primarily centralized; limited cryptographic, shared-ledger guarantees\n"
        "3) SAP ATTP / SAP ICH for Life Sciences\n"
        "   Strengths: Deep ERP integration; enterprise-grade operations\n"
        "   Weaknesses: Vendor lock-in; cross-company transparency depends on integrations\n\n"
        "PharmaChain differentiation: cryptographic provenance (hash chain + Merkle root + RSA signatures) + role-based dashboards with verification endpoints.",
    )

    # Slide 5 — Decision Matrix (Optional)
    s5 = prs.slides[4]
    s5.placeholders[13].insert_picture(str(decision_img))

    # Slide 6 — Opportunity Identification
    s6 = prs.slides[5]
    _set_text(s6.placeholders[14], "PharmaChain")
    _set_text(
        s6.placeholders[16],
        "PharmaChain is a blockchain-inspired pharmaceutical supply tracker that records medicine batch events (creation, transfer, dispensing) with cryptographic proofs.\n"
        "It supports companies, hospitals, and customers to track and verify batch authenticity using hash chaining, Merkle-root integrity checks, and RSA digital signatures.",
    )
    _set_text(
        s6.placeholders[17],
        "• High-risk domain where tampering and counterfeits directly harm patients\n"
        "• Multiple organizations must coordinate but don’t fully trust each other\n"
        "• Immutable, signed events improve auditability and accountability\n"
        "• Proof storage (hashes/Merkle roots) provides integrity guarantees with low data exposure\n"
        "• Fits a permissioned governance model aligned with real-world compliance needs",
    )

    # Slide 7 — Statements
    s7 = prs.slides[6]
    _set_text(
        s7.placeholders[17],
        "Patients, pharmacies, and hospitals lack a single trusted way to verify medicine authenticity and track batch custody across the full supply chain, leading to counterfeit risk, recall delays, and high reconciliation overhead.",
    )
    _set_text(
        s7.placeholders[18],
        "PharmaChain provides a tamper-evident, verifiable provenance trail for each medicine batch so stakeholders can detect anomalies faster, improve recall precision, reduce disputes, and give customers confidence in what they purchased.",
    )

    # Slide 8 — Value + Positioning statement
    s8 = prs.slides[7]
    _set_text(
        s8.placeholders[17],
        "• Verifiable batch provenance across organizations\n"
        "• Strong authenticity guarantees via RSA signatures\n"
        "• Efficient integrity checks via Merkle proofs\n"
        "• Faster, more targeted recalls and compliance reporting\n"
        "• Customer-facing verification (track/verify endpoints + dashboard)",
    )
    _set_text(
        s8.placeholders[18],
        "For manufacturers, hospitals, pharmacies, and customers who need trustworthy medicine provenance, PharmaChain is a permissioned traceability platform that records signed batch events and integrity proofs, enabling instant verification and audit-ready transparency across the supply chain.",
    )

    # Slide 9 — Strategic positioning
    s9 = prs.slides[8]
    _set_text(
        s9.placeholders[17],
        "PharmaChain positions an organization as a trusted, compliance-forward supply chain partner by providing verifiable provenance and faster response to recalls and counterfeit incidents.\n\n"
        "Sustaining the positioning:\n"
        "• Start with pilots (1 manufacturer + 1 distributor + 2 pharmacies/hospitals)\n"
        "• Integrate with existing workflows (APIs, CSV import/export, ERP connectors)\n"
        "• Establish governance and permissioning for who can write/verify events\n"
        "• Use standards (serialization/GS1 identifiers) to reduce friction and scale",
    )

    # Slide 10 — Execution (Business Model Canvas)
    s10 = prs.slides[9]
    s10.placeholders[13].insert_picture(str(bmc_img))

    # Slide 11 — Business model description
    s11 = prs.slides[10]
    _set_text(
        s11.placeholders[17],
        "Key activities: participant onboarding, ERP/workflow integrations, blockchain/node operations, smart contract + API maintenance, and compliance reporting.\n"
        "Cost structure: engineering, cloud + node hosting, security audits, support/customer success, and governance overhead.\n"
        "Revenue streams: subscription pricing per participant (manufacturer/distributor/pharmacy), tiered by transaction volume, plus premium analytics/compliance reports and integration services.",
    )

    # Slide 12 — Team + Partners
    s12 = prs.slides[11]
    _set_text(
        s12.placeholders[17],
        "Team needs:\n"
        "• Blockchain/cryptography engineer (signatures, integrity proofs)\n"
        "• Backend engineer (Rust/Axum, SQLx)\n"
        "• Frontend engineer (Next.js, dashboards)\n"
        "• Security engineer (key management, audits)\n"
        "• Product/compliance lead (pharma workflows, regulation)",
    )
    _set_text(
        s12.placeholders[18],
        "Partners & allies:\n"
        "• Manufacturers and distributors for pilot deployments\n"
        "• Pharmacies/hospitals for last-mile verification\n"
        "• Regulators/industry bodies for traceability alignment\n"
        "• GS1/serialization ecosystem + integration vendors\n"
        "• Security/audit firms for assurance",
    )

    # Slide 13 — Operational considerations
    s13 = prs.slides[12]
    _set_text(
        s13.placeholders[17],
        "Operational considerations / boundaries:\n"
        "• Adoption risk: incentives and low-friction integration are critical\n"
        "• Data quality: validate identifiers, timestamps, and role permissions\n"
        "• Privacy: keep sensitive fields off-chain; store only hashes/proofs where needed\n"
        "• Governance: permissioned membership, key management, and audit policies\n"
        "• Performance: batch writes, async processing, and scalable verification endpoints\n\n"
        "Mitigations: phased rollout, clear governance model, API-first integrations, and security reviews.",
    )

    # Slide 14 — Closing
    s14 = prs.slides[13]
    _set_text(s14.placeholders[0], "Thank you")
    _set_text(s14.placeholders[1], "Lohitaksha Patary")
    _set_text(s14.placeholders[11], "Undergraduate Student, Amrita Vishwa Vidyapeetham, Bangalore")

    prs.save(str(OUT_PPTX))
    print(f"Wrote: {OUT_PPTX}")


if __name__ == "__main__":
    main()
